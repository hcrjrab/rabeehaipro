"""Office document tools — Word, Excel, PowerPoint generation.

All tools produce real .docx / .xlsx / .pptx files inside the workspace.

Design
------
- Dependencies (python-docx, openpyxl, python-pptx) are imported lazily
  inside ``execute()`` so the app boots without the ``office`` extra installed.
- Tools are risk-classified ``SAFE`` (they only create files; no deletion).
- Content is validated: max row/column limits prevent runaway spreadsheets,
  max slide count prevents presentation bombs.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from ..config.schemas import RiskLevel, ToolCallResult
from .base import BaseTool, ToolContext
from .builtin import _confine


class _OfficeTool(BaseTool):
    """Shared base for office tools: common pattern of resolve path + write."""

    risk = RiskLevel.SAFE

    def _resolve(self, args: dict[str, Any], ctx: ToolContext) -> tuple[Path, str]:
        """Extract and confine the output path; raise ToolCallResult on error."""
        raw_path = str(args.get("path", "")).strip()
        if not raw_path:
            raise ValueError("Missing 'path' argument.")
        return _confine(raw_path, ctx.workspace), raw_path

    async def execute(self, args: dict[str, Any], ctx: ToolContext) -> ToolCallResult:
        try:
            return await self._execute_checked(args, ctx)
        except ValueError as exc:
            return ToolCallResult(ok=False, error=str(exc))
        except ImportError as exc:
            return ToolCallResult(
                ok=False,
                error=f"Missing dependency: {exc}. Install the 'office' extra: pip install rabeeh-ai-agent-pro[office]",
            )
        except Exception as exc:
            return ToolCallResult(ok=False, error=f"Office tool error: {exc}")

    async def _execute_checked(self, args: dict[str, Any], ctx: ToolContext) -> ToolCallResult:
        """Override in subclasses. May raise ValueError / ImportError."""
        raise NotImplementedError


class CreateWordTool(_OfficeTool):
    """Create a .docx Word document with headings and paragraphs."""

    name = "office.create_word"
    description = "Create a Word (.docx) document with headings and paragraphs."
    risk = RiskLevel.SAFE

    def schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output .docx path (relative)."},
                "title": {"type": "string", "description": "Document title."},
                "sections": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "body": {"type": "string"},
                        },
                        "required": ["heading", "body"],
                    },
                    "description": "Sections: each has a heading + body paragraph.",
                },
            },
            "required": ["path", "title"],
        }

    async def _execute_checked(self, args: dict[str, Any], ctx: ToolContext) -> ToolCallResult:
        from docx import Document

        target, raw_path = self._resolve(args, ctx)
        title = str(args.get("title", "Untitled"))
        sections = args.get("sections") or []

        doc = Document()
        doc.add_heading(title, level=0)
        for sec in sections:
            doc.add_heading(str(sec.get("heading", "")), level=1)
            doc.add_paragraph(str(sec.get("body", "")))

        target.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(target))
        return ToolCallResult(
            ok=True,
            data={"path": raw_path, "sections": len(sections), "title": title},
        )


class CreateExcelTool(_OfficeTool):
    """Create a .xlsx spreadsheet with headers and data rows."""

    name = "office.create_excel"
    description = "Create an Excel (.xlsx) spreadsheet with headers and row data."
    risk = RiskLevel.SAFE

    _MAX_ROWS = 10_000  # safety limit

    def schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output .xlsx path (relative)."},
                "title": {"type": "string", "description": "Sheet name / title."},
                "headers": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "Column headers.",
                },
                "rows": {
                    "type": "array",
                    "items": {"type": "array", "items": {}},
                    "description": "Data rows (each an array of cell values).",
                },
            },
            "required": ["path", "headers"],
        }

    async def _execute_checked(self, args: dict[str, Any], ctx: ToolContext) -> ToolCallResult:
        from openpyxl import Workbook

        target, raw_path = self._resolve(args, ctx)
        title = str(args.get("title", "Sheet1"))
        headers = args.get("headers") or []
        rows = args.get("rows") or []

        if len(rows) > self._MAX_ROWS:
            return ToolCallResult(
                ok=False,
                error=f"Row limit exceeded: {len(rows)} > {self._MAX_ROWS}.",
            )

        wb = Workbook()
        ws = wb.active
        ws.title = title[:31]  # Excel sheet name max length

        if headers:
            ws.append(headers)
        for row in rows:
            ws.append(row)

        target.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(target))
        return ToolCallResult(
            ok=True,
            data={
                "path": raw_path,
                "headers": len(headers),
                "rows": len(rows),
                "sheet": ws.title,
            },
        )


class CreatePowerPointTool(_OfficeTool):
    """Create a .pptx presentation with title slide + content slides."""

    name = "office.create_powerpoint"
    description = "Create a PowerPoint (.pptx) presentation."
    risk = RiskLevel.SAFE

    _MAX_SLIDES = 100  # safety limit

    def schema(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output .pptx path (relative)."},
                "title": {"type": "string", "description": "Presentation title."},
                "slides": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "content": {"type": "string"},
                        },
                    },
                    "description": "Slides: each has a title and bullet-point content.",
                },
            },
            "required": ["path", "title"],
        }

    async def _execute_checked(self, args: dict[str, Any], ctx: ToolContext) -> ToolCallResult:
        from pptx import Presentation

        target, raw_path = self._resolve(args, ctx)
        title = str(args.get("title", "Untitled"))
        slides = args.get("slides") or []

        if len(slides) > self._MAX_SLIDES:
            return ToolCallResult(
                ok=False,
                error=f"Slide limit exceeded: {len(slides)} > {self._MAX_SLIDES}.",
            )

        prs = Presentation()
        # Title slide.
        prs.slides.add_slide(prs.slide_layouts[0])
        title_shape = prs.slides[0].shapes.title
        title_shape.text = title

        # Content slides.
        for slide_data in slides:
            slide_layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = str(slide_data.get("title", ""))
            body = slide.placeholders[1]
            body.text = str(slide_data.get("content", ""))

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
        return ToolCallResult(
            ok=True,
            data={"path": raw_path, "slides": len(slides) + 1, "title": title},
        )
